import chromadb
from chromadb.config import Settings
import os
try:
    from PyPDF2 import PdfReader
except ImportError:
    from pypdf import PdfReader
from docx import Document
from pptx import Presentation

class VectorService:
    def __init__(self):
        # Use the new ChromaDB client configuration with telemetry disabled
        self.client = chromadb.PersistentClient(
            path="./chroma_db",
            settings=Settings(
                anonymized_telemetry=False,
                allow_reset=True
            )
        )
        self.collection = self.client.get_or_create_collection("files")
    
    def extract_text(self, file_path, file_type):
        """Extract text from various file types"""
        try:
            if file_type == 'pdf':
                reader = PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
            
            elif file_type in ['doc', 'docx']:
                doc = Document(file_path)
                return "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            elif file_type in ['ppt', 'pptx']:
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            
            elif file_type == 'txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            
            else:
                # For media files, return filename as placeholder
                return f"Media file: {os.path.basename(file_path)}"
        
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return ""
    
    def add_file(self, file_path, file_id):
        """Add file to vector database"""
        from models import File
        file = File.query.get(file_id)
        
        text = self.extract_text(file_path, file.file_type)
        
        if text.strip():
            # Split into chunks (simple chunking - can be improved)
            chunks = self._chunk_text(text)
            
            # Add to ChromaDB
            for idx, chunk in enumerate(chunks):
                self.collection.add(
                    documents=[chunk],
                    metadatas=[{"file_id": file_id, "chunk": idx, "filename": file.filename}],
                    ids=[f"file_{file_id}_chunk_{idx}"]
                )
        
        return f"file_{file_id}"
    
    def _chunk_text(self, text, chunk_size=1000, overlap=200):
        """Split text into overlapping chunks"""
        chunks = []
        start = 0
        text_len = len(text)
        
        while start < text_len:
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start += (chunk_size - overlap)
        
        return chunks
    
    def get_files_context(self, file_ids):
        """Get overview context for files"""
        if not file_ids:
            return "No content available"
        
        try:
            results = self.collection.get(
                where={"file_id": {"$in": file_ids}},
                limit=10
            )
            
            if results and results['documents']:
                return "\n\n".join(results['documents'][:5])  # Return first 5 chunks as context
        except Exception as e:
            print(f"Error getting files context: {e}")
        
        return "No content available"
    
    def query_relevant_content(self, queries, file_ids, n_results=5):
        """Query vector DB for relevant content"""
        if not queries or not file_ids:
            return ""
        
        # Combine queries into one search
        combined_query = " ".join(queries) if isinstance(queries, list) else queries
        
        try:
            results = self.collection.query(
                query_texts=[combined_query],
                where={"file_id": {"$in": file_ids}},
                n_results=n_results
            )
            
            if results and results['documents']:
                return "\n\n".join(results['documents'][0])
            return ""
        except Exception as e:
            print(f"Query error: {e}")
            return self.get_files_context(file_ids)
